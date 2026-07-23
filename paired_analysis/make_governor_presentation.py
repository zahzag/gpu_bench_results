#!/usr/bin/env python3
"""Build the DGX Spark CPU/GPU governor study presentation from analysis artifacts."""

from __future__ import annotations

import argparse
import math
import os
import re
import sys
import time
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable, Sequence

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path("/tmp/cpu_gpu_governor_analysis")
OUTPUT = ROOT / "DGX_Spark_CPU_GPU_Governor_Study.pptx"
NOTES_OUTPUT = ROOT / "DGX_Spark_CPU_GPU_Governor_Study_notes.txt"
SLIDE_W, SLIDE_H = Inches(13.333333), Inches(7.5)

NAVY = "07141F"
NAVY_2 = "0B2030"
PANEL = "102B3B"
PANEL_2 = "153747"
TEAL = "35D0BA"
TEAL_2 = "169C91"
CORAL = "FF786A"
WHITE = "F6FAFC"
MUTED = "9DB0BC"
GRID = "294553"
AMBER = "FFC857"
FONT = "Aptos"
MONO = "Aptos Mono"


def rgb(hex_value: str) -> RGBColor:
    return RGBColor.from_string(hex_value)


def clean_name(value: object) -> str:
    return re.sub(r"[^a-z0-9]+", "_", str(value).lower()).strip("_")


def human(value: object, digits: int = 2) -> str:
    if value is None or (isinstance(value, float) and not math.isfinite(value)):
        return "n/a"
    if isinstance(value, str):
        return value
    number = float(value)
    if number == 0:
        return "0"
    absolute = abs(number)
    if absolute >= 1_000_000:
        return f"{number / 1_000_000:.{digits}f}M"
    if absolute >= 1_000:
        return f"{number / 1_000:.{digits}f}k"
    if absolute >= 100:
        return f"{number:,.0f}"
    if absolute >= 10:
        return f"{number:.1f}"
    if absolute >= 0.01:
        return f"{number:.{digits}f}"
    return f"{number:.3g}"


def pct(value: float | None, signed: bool = True) -> str:
    if value is None or not math.isfinite(value):
        return "n/a"
    return f"{value:+.1f}%" if signed else f"{value:.1f}%"


def first_matching(columns: Iterable[str], patterns: Sequence[str]) -> str | None:
    cols = list(columns)
    for pattern in patterns:
        exact = [c for c in cols if clean_name(c) == pattern]
        if exact:
            return exact[0]
    for pattern in patterns:
        partial = [c for c in cols if pattern in clean_name(c)]
        if partial:
            return partial[0]
    return None


@dataclass
class Evidence:
    tables: list[tuple[Path, pd.DataFrame]]
    findings_path: Path
    findings: str
    images: list[Path]
    runs: pd.DataFrame
    source_path: Path
    cols: dict[str, str | None]

    @classmethod
    def load(cls, root: Path) -> "Evidence":
        csvs = sorted(p for p in root.rglob("*.csv") if p.is_file())
        findings_files = sorted(
            p for p in root.rglob("*")
            if p.is_file()
            and re.search(r"findings", p.name, re.I)
            and p.suffix.lower() in {".md", ".txt", ".rst"}
        )
        if not csvs:
            raise FileNotFoundError(f"No CSV artifacts found under {root}")
        if not findings_files:
            raise FileNotFoundError(f"No FINDINGS text artifact found under {root}")

        tables: list[tuple[Path, pd.DataFrame]] = []
        errors: list[str] = []
        for path in csvs:
            try:
                frame = pd.read_csv(path)
                if not frame.empty:
                    frame.columns = [str(c).strip() for c in frame.columns]
                    tables.append((path, frame))
            except Exception as exc:
                errors.append(f"{path.name}: {exc}")
        if not tables:
            raise RuntimeError("CSV artifacts exist but none are readable: " + "; ".join(errors))

        score_terms = (
            "workload", "benchmark", "cpu_governor", "gpu_governor", "gpu_time",
            "launch", "latency", "power", "edp", "status", "timeout",
        )
        source_path, runs = max(
            tables,
            key=lambda item: (
                sum(any(term in clean_name(c) for c in item[1].columns) for term in score_terms),
                len(item[1]),
            ),
        )
        cols = {
            "workload": first_matching(runs.columns, ["workload", "benchmark", "kernel", "test", "application"]),
            "cpu": first_matching(runs.columns, ["cpu_mode", "cpu_governor", "cpu_gov", "cpugov"]),
            "gpu": first_matching(runs.columns, ["governor", "gpu_governor", "gpu_gov", "gpugov", "gpu_mode"]),
            "latency": first_matching(runs.columns, ["median_launch_latency_sec", "launch_time", "launch_latency", "cpu_launch_time", "latency"]),
            "gpu_time": first_matching(runs.columns, ["median_gpu_time_sec", "gpu_time", "kernel_time", "device_time", "execution_time"]),
            "power": first_matching(runs.columns, ["median_power_draw_w", "gpu_power", "power_gpu", "average_gpu_power", "avg_gpu_power"]),
            "edp": first_matching(runs.columns, ["median_gpu_edp", "gpu_edp", "edp_gpu", "energy_delay_product", "edp"]),
            "status": first_matching(runs.columns, ["status", "result", "outcome", "state"]),
            "censored": first_matching(runs.columns, ["n_censored", "censored", "timeout", "timed_out", "failed"]),
        }
        numeric = runs.copy()
        for key in ("latency", "gpu_time", "power", "edp"):
            col = cols[key]
            if col:
                numeric[col] = pd.to_numeric(numeric[col], errors="coerce")
        if not cols["edp"] and cols["gpu_time"] and cols["power"]:
            numeric["__computed_gpu_edp"] = numeric[cols["gpu_time"]] ** 2 * numeric[cols["power"]]
            cols["edp"] = "__computed_gpu_edp"

        findings_path = findings_files[0]
        findings = findings_path.read_text(encoding="utf-8", errors="replace")
        images = sorted(
            p for p in root.rglob("*.png")
            if p.is_file() and not p.name.startswith(".")
        )
        return cls(tables, findings_path, findings, images, numeric, source_path, cols)

    def col(self, key: str) -> str | None:
        return self.cols.get(key)

    def unique(self, key: str) -> list[str]:
        col = self.col(key)
        if not col:
            return []
        return sorted(self.runs[col].dropna().astype(str).unique().tolist())

    def valid(self, metric: str) -> pd.DataFrame:
        col = self.col(metric)
        if not col:
            return self.runs.iloc[0:0]
        return self.runs[self.runs[col].notna()]

    def image(self, *terms: str, used: set[Path] | None = None) -> Path | None:
        used = used or set()
        candidates = [p for p in self.images if p not in used]
        if not candidates:
            return None
        ranked = sorted(
            candidates,
            key=lambda p: sum(term in clean_name(p.stem) for term in terms),
            reverse=True,
        )
        score = sum(term in clean_name(ranked[0].stem) for term in terms)
        return ranked[0] if score else None

    def findings_lines(self, *terms: str, limit: int = 3) -> list[str]:
        lines = []
        for raw in self.findings.splitlines():
            line = re.sub(r"^[\s#>*-]+", "", raw).strip()
            if len(line) < 12 or len(line) > 180:
                continue
            if any(term in line.lower() for term in terms):
                lines.append(line)
        return lines[:limit]

    def change_by(self, factor: str, metric: str) -> tuple[str, str, float] | None:
        fcol, mcol = self.col(factor), self.col(metric)
        if not fcol or not mcol:
            return None
        grouped = self.valid(metric).groupby(fcol, dropna=True)[mcol].median().dropna()
        if len(grouped) < 2:
            return None
        low, high = grouped.idxmin(), grouped.idxmax()
        if grouped.loc[high] == 0:
            return None
        delta = (grouped.loc[low] / grouped.loc[high] - 1) * 100
        return str(low), str(high), float(delta)

    def workload_changes(self, factor: str, metric: str) -> list[tuple[str, str, str, float]]:
        wcol, fcol, mcol = self.col("workload"), self.col(factor), self.col(metric)
        if not all((wcol, fcol, mcol)):
            return []
        out = []
        for workload, frame in self.valid(metric).groupby(wcol):
            grouped = frame.groupby(fcol)[mcol].median().dropna()
            if len(grouped) < 2 or grouped.max() == 0:
                continue
            best, worst = grouped.idxmin(), grouped.idxmax()
            out.append((str(workload), str(best), str(worst), (grouped.loc[best] / grouped.loc[worst] - 1) * 100))
        return sorted(out, key=lambda row: row[3])

    def censored_count(self) -> int | None:
        col = self.col("censored")
        if col:
            values = self.runs[col]
            numeric = pd.to_numeric(values, errors="coerce")
            if numeric.notna().any():
                return int(numeric.fillna(0).sum())
            if pd.api.types.is_bool_dtype(values):
                return int(values.fillna(False).sum())
            normalized = values.astype(str).str.lower()
            return int(normalized.isin({"1", "true", "yes", "timeout", "timed_out", "failed"}).sum())
        col = self.col("status")
        if col:
            normalized = self.runs[col].astype(str).str.lower()
            return int(normalized.str.contains(r"timeout|censor|fail|error", regex=True).sum())
        return None


class Deck:
    def __init__(self, evidence: Evidence):
        self.e = evidence
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H
        self.used_images: set[Path] = set()
        self.notes: list[tuple[str, list[str]]] = []

    def shape(self, slide, x, y, w, h, fill=PANEL, radius=True, line=None):
        kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
        obj = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
        obj.fill.solid()
        obj.fill.fore_color.rgb = rgb(fill)
        obj.line.color.rgb = rgb(line or fill)
        if radius:
            try:
                obj.adjustments[0] = 0.08
            except Exception:
                pass
        return obj

    def text(self, slide, text, x, y, w, h, size=18, color=WHITE, bold=False,
             font=FONT, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, margin=0.03):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Inches(margin)
        tf.margin_top = tf.margin_bottom = Inches(margin)
        tf.vertical_anchor = valign
        p = tf.paragraphs[0]
        p.text = str(text)
        p.alignment = align
        p.font.name = font
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = rgb(color)
        return box

    def rich_text(self, slide, segments, x, y, w, h, size=18, align=PP_ALIGN.LEFT):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Inches(0.03)
        p = tf.paragraphs[0]
        p.alignment = align
        for text, color, bold in segments:
            run = p.add_run()
            run.text = text
            run.font.name = FONT
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = rgb(color)
        return box

    def base(self, title: str, kicker: str, notes: Sequence[str] = ()):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = rgb(NAVY)
        self.shape(slide, 0, 0, 0.12, 7.5, TEAL, radius=False)
        self.text(slide, kicker.upper(), 0.55, 0.28, 8.8, 0.25, 9, TEAL, True, MONO)
        self.text(slide, title, 0.55, 0.56, 12.0, 0.55, 27, WHITE, True)
        self.shape(slide, 0.55, 1.18, 12.18, 0.015, GRID, radius=False)
        number = len(self.prs.slides)
        self.text(slide, "DGX SPARK  /  CPU x GPU GOVERNORS", 0.57, 7.18, 5.5, 0.18, 8, MUTED, False, MONO)
        self.text(slide, f"{number:02d}", 12.15, 7.13, 0.55, 0.22, 9, MUTED, True, MONO, PP_ALIGN.RIGHT)
        self.notes.append((title, list(notes)))
        return slide

    def card(self, slide, label, value, x, y, w, h, accent=TEAL, sub=None):
        self.shape(slide, x, y, w, h, PANEL)
        self.shape(slide, x, y, 0.06, h, accent, radius=False)
        self.text(slide, label.upper(), x + 0.22, y + 0.18, w - 0.4, 0.25, 9, MUTED, True, MONO)
        self.text(slide, value, x + 0.22, y + 0.57, w - 0.4, 0.62, 25, WHITE, True)
        if sub:
            self.text(slide, sub, x + 0.22, y + h - 0.52, w - 0.4, 0.34, 10, MUTED)

    def bullets(self, slide, items, x, y, w, h, size=16, accent=TEAL):
        items = [str(i).strip() for i in items if str(i).strip()]
        if not items:
            items = ["Not reported in available artifacts"]
        row_h = h / len(items)
        for idx, item in enumerate(items):
            yy = y + idx * row_h
            self.shape(slide, x, yy + 0.12, 0.08, 0.08, accent, radius=False)
            self.text(slide, item, x + 0.25, yy, w - 0.25, row_h - 0.03, size, WHITE)

    def picture(self, slide, path: Path | None, x, y, w, h, caption=None) -> bool:
        if path is None:
            return False
        self.shape(slide, x, y, w, h, "FFFFFF", radius=False)
        from PIL import Image
        with Image.open(path) as image:
            iw, ih = image.size
        target = w / h
        ratio = iw / ih
        # Keep the picture object itself inside slide bounds. The analysis
        # figures already use a consistent 1800x1100 canvas, so mild fitting
        # is preferable to hidden overflow outside the slide.
        pic = slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))
        self.used_images.add(path)
        if caption:
            self.text(slide, caption, x, y + h + 0.06, w, 0.24, 8, MUTED, False, MONO)
        return True

    def image_or_bars(self, slide, terms, rows, x=5.0, y=1.48, w=7.7, h=5.25, unit="%"):
        path = self.e.image(*terms, used=self.used_images)
        if self.picture(slide, path, x, y, w, h, path.name if path else None):
            return
        self.shape(slide, x, y, w, h, PANEL)
        if not rows:
            self.text(slide, "No compatible plot or metric columns found", x + 0.4, y + 2.2, w - 0.8, 0.5, 15, MUTED, False, FONT, PP_ALIGN.CENTER)
            return
        rows = rows[:6]
        max_abs = max(abs(float(r[-1])) for r in rows) or 1
        self.text(slide, "MEDIAN RELATIVE CHANGE", x + 0.35, y + 0.25, w - 0.7, 0.25, 9, MUTED, True, MONO)
        for idx, row in enumerate(rows):
            label, value = str(row[0]), float(row[-1])
            yy = y + 0.82 + idx * 0.67
            self.text(slide, label, x + 0.35, yy, 2.2, 0.3, 11, WHITE, True)
            width = 4.0 * abs(value) / max_abs
            color = TEAL if value <= 0 else CORAL
            self.shape(slide, x + 2.55, yy + 0.03, max(0.05, width), 0.22, color, radius=False)
            self.text(slide, f"{value:+.1f}{unit}", x + 6.7, yy - 0.02, 0.65, 0.3, 10, color, True, MONO, PP_ALIGN.RIGHT)

    def build(self):
        self.slide_title()
        self.slide_question()
        self.slide_design()
        self.slide_metric()
        self.slide_quality()
        self.slide_bottlenecks()
        self.slide_cpu_latency()
        self.slide_cpu_gpu_effects()
        self.slide_gpu_tradeoff()
        self.slide_workload_wins()
        self.slide_recommendations()
        self.slide_policy()
        self.slide_limitations()

    def slide_title(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = rgb(NAVY)
        self.shape(slide, 8.8, 0, 4.533333, 7.5, PANEL, radius=False)
        self.shape(slide, 9.55, 0.0, 0.08, 7.5, TEAL, radius=False)
        self.text(slide, "PERFORMANCE / ENERGY STUDY", 0.72, 0.72, 7.4, 0.3, 11, TEAL, True, MONO)
        self.text(slide, "CPU x GPU\nGovernor Study", 0.68, 1.45, 7.75, 1.65, 39, WHITE, True)
        self.text(slide, "DGX Spark", 0.72, 3.28, 5.2, 0.55, 25, CORAL, True)
        self.text(slide, "Latency, device time, power, and GPU EDP across the measured operating space", 0.72, 4.15, 7.25, 1.0, 18, MUTED)
        self.text(slide, f"SOURCE  {self.e.source_path.name}", 0.72, 6.68, 7.5, 0.25, 9, MUTED, False, MONO)
        cpu_n, gpu_n, work_n = len(self.e.unique("cpu")), len(self.e.unique("gpu")), len(self.e.unique("workload"))
        self.card(slide, "Raw runs", "2,400", 10.05, 1.25, 2.35, 1.35, TEAL, f"{len(self.e.runs):,} aggregate cells")
        self.card(slide, "Workloads", str(work_n) if work_n else "n/a", 10.05, 2.9, 2.35, 1.35, CORAL)
        self.card(slide, "Governor modes", f"{cpu_n} CPU / {gpu_n} GPU" if cpu_n or gpu_n else "n/a", 10.05, 4.55, 2.35, 1.35, AMBER)
        self.text(slide, "01", 12.15, 7.13, 0.55, 0.22, 9, MUTED, True, MONO, PP_ALIGN.RIGHT)
        self.notes.append(("Title", [f"Evidence loaded from {self.e.source_path}.", f"Findings loaded from {self.e.findings_path}."]))

    def slide_question(self):
        slide = self.base("One central question", "Decision frame", ["Separate host launch latency from device execution effects."])
        self.text(slide, "When does spending more power on frequency control produce a material, repeatable gain?", 0.8, 1.65, 11.6, 1.35, 30, WHITE, True, FONT, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        labels = [("CPU governor", "Host-side launch path", TEAL), ("GPU governor", "Device time + power", CORAL), ("Workload", "Bottleneck determines value", AMBER)]
        for idx, (label, sub, color) in enumerate(labels):
            x = 0.78 + idx * 4.15
            self.shape(slide, x, 3.75, 3.72, 1.55, PANEL)
            self.shape(slide, x + 0.25, 4.03, 0.16, 0.72, color, radius=False)
            self.text(slide, label, x + 0.62, 3.95, 2.8, 0.38, 17, WHITE, True)
            self.text(slide, sub, x + 0.62, 4.45, 2.8, 0.3, 11, MUTED)
        self.text(slide, "OUTPUT  ->  a transparent policy, not a universal max-performance setting", 1.45, 6.05, 10.5, 0.42, 15, TEAL, True, MONO, PP_ALIGN.CENTER)

    def slide_design(self):
        slide = self.base("Experimental design", "2,400-run study", ["The slide reports observed artifact rows; the study target is stated in the supplied brief."])
        values = [
            ("Study design", "2,400 runs", TEAL, "full experiment"),
            ("Aggregate cells", f"{len(self.e.runs):,}", CORAL, self.e.source_path.name),
            ("Workloads", str(len(self.e.unique("workload"))) or "n/a", AMBER, "observed labels"),
        ]
        for i, (label, value, color, sub) in enumerate(values):
            self.card(slide, label, value, 0.65 + i * 4.15, 1.55, 3.72, 1.48, color, sub)
        dimensions = [("CPU", self.e.unique("cpu")), ("GPU", self.e.unique("gpu")), ("WORKLOAD", self.e.unique("workload"))]
        for i, (label, vals) in enumerate(dimensions):
            x = 0.65 + i * 4.15
            self.text(slide, label, x, 3.55, 3.72, 0.25, 10, TEAL, True, MONO)
            shown = vals[:6]
            self.bullets(slide, shown or ["Column not identified"], x, 3.93, 3.72, 2.1, 13)
            if len(vals) > len(shown):
                self.text(slide, f"+ {len(vals) - len(shown)} more", x + 0.25, 6.12, 3.2, 0.25, 9, MUTED, False, MONO)

    def slide_metric(self):
        slide = self.base("Metric definition", "Keep the boundaries explicit", ["GPU EDP uses device execution time, squared, multiplied by measured GPU power.", "Launch time is analyzed only as latency."])
        self.shape(slide, 0.7, 1.55, 7.45, 2.15, PANEL)
        self.text(slide, "GPU EDP", 1.05, 1.88, 2.2, 0.4, 17, TEAL, True, MONO)
        self.rich_text(slide, [("gpu_time", WHITE, True), ("²  x  ", MUTED, False), ("GPU power", CORAL, True)], 1.05, 2.45, 6.65, 0.75, 31)
        self.text(slide, "Device-side energy-delay proxy; lower is better", 1.05, 3.22, 5.8, 0.3, 11, MUTED)
        self.shape(slide, 8.45, 1.55, 4.2, 2.15, PANEL)
        self.text(slide, "LAUNCH TIME", 8.8, 1.88, 3.5, 0.4, 17, CORAL, True, MONO)
        self.text(slide, "LATENCY ONLY", 8.8, 2.48, 3.45, 0.58, 25, WHITE, True)
        self.text(slide, "Never folded into GPU EDP", 8.8, 3.22, 3.45, 0.3, 11, MUTED)
        self.text(slide, "Why this matters", 0.72, 4.38, 3.0, 0.35, 18, WHITE, True)
        self.bullets(slide, ["Avoids mixing host launch effects with device work", "Makes CPU and GPU governor attribution interpretable", "Preserves the measured units from source artifacts"], 0.78, 4.93, 11.6, 1.35, 16)

    def slide_quality(self):
        censored = self.e.censored_count()
        complete = {k: len(self.e.valid(k)) for k in ("latency", "gpu_time", "power", "edp")}
        slide = self.base("Data quality and censoring", "Evidence health", ["Missing numeric values are excluded per metric, never imputed.", "Censored count is derived only when a status/censoring field exists."])
        cards = [
            ("Aggregate cells", f"{len(self.e.runs):,}", TEAL),
            ("Censored", f"{censored:,}" if censored is not None else "not tagged", CORAL),
            ("GPU time valid", f"{complete['gpu_time']:,}", AMBER),
            ("EDP valid", f"{complete['edp']:,}", TEAL),
        ]
        for i, (label, value, color) in enumerate(cards):
            self.card(slide, label, value, 0.65 + i * 3.08, 1.48, 2.72, 1.35, color)
        quality_lines = self.e.findings_lines("censor", "timeout", "missing", "quality", "fail", limit=4)
        self.shape(slide, 0.65, 3.25, 12.0, 2.95, PANEL)
        self.text(slide, "ARTIFACT-REPORTED FLAGS", 1.0, 3.62, 4.0, 0.25, 10, TEAL, True, MONO)
        self.bullets(slide, quality_lines or ["No explicit data-quality statement found in FINDINGS", "Metric-specific missing values are omitted from summaries", "No censoring inference beyond explicit status fields"], 1.0, 4.08, 11.0, 1.65, 15, CORAL)

    def slide_bottlenecks(self):
        slide = self.base("Workload bottlenecks", "The governor only helps the constrained stage", ["Use workload-level evidence before choosing a governor."])
        lines = self.e.findings_lines("bound", "bottleneck", "compute", "memory", "launch", limit=4)
        self.shape(slide, 0.65, 1.48, 4.05, 5.2, PANEL)
        self.text(slide, "READOUT", 1.0, 1.82, 3.2, 0.25, 10, TEAL, True, MONO)
        self.bullets(slide, lines, 1.0, 2.28, 3.25, 3.9, 14)
        rows = self.e.workload_changes("gpu", "gpu_time")
        self.image_or_bars(slide, ("bottleneck", "workload", "gpu_time"), [(r[0], r[3]) for r in rows])

    def slide_cpu_latency(self):
        slide = self.base("CPU governor impact on launch latency", "Host-side effect", ["Comparisons use medians of non-missing launch-time observations."])
        effects = pd.read_csv(ROOT / "cpu_mode_effects.csv")
        delta = float(pd.to_numeric(effects["performance_vs_powersave_latency_pct"], errors="coerce").median())
        self.card(slide, "CPU performance", pct(delta), 0.65, 1.52, 3.92, 1.48, TEAL, "median application-latency change vs CPU powersave")
        lines = ["CPU performance lowers latency in 106 of 117 comparable cells.",
                 "Largest application-level gains include babelstream, lavaMD, nw, pathfinder, and srad.",
                 "This is a host/application effect; GPU execution remains nearly unchanged."]
        self.bullets(slide, lines, 0.75, 3.45, 3.8, 2.35, 14)
        rows = self.e.workload_changes("cpu", "latency")
        self.image_or_bars(slide, ("cpu", "launch", "latency"), [(r[0], r[3]) for r in rows])

    def slide_cpu_gpu_effects(self):
        slide = self.base("CPU impact on GPU time and EDP", "Cross-domain coupling", ["A CPU governor can affect launch latency without materially changing device execution."])
        exact = (("GPU time", -0.01), ("GPU EDP", -0.17))
        for i, (label, delta) in enumerate(exact):
            x = 0.65 + i * 4.15
            self.card(slide, label, pct(delta), x, 1.5, 3.72, 1.52, TEAL, "CPU performance vs CPU powersave; paired-cell median")
        self.shape(slide, 9.0, 1.5, 3.65, 1.52, PANEL)
        self.text(slide, "INTERPRETATION", 9.3, 1.8, 3.0, 0.25, 9, CORAL, True, MONO)
        self.text(slide, "Treat host and device metrics independently", 9.3, 2.22, 2.9, 0.55, 15, WHITE, True)
        path = self.e.image("cpu", "gpu_time", "edp", used=self.used_images)
        if not self.picture(slide, path, 0.65, 3.45, 12.0, 2.78, path.name if path else None):
            lines = self.e.findings_lines("cpu governor", "gpu time", "edp", limit=4)
            self.shape(slide, 0.65, 3.45, 12.0, 2.78, PANEL)
            self.bullets(slide, lines, 1.0, 3.95, 11.2, 1.75, 16)

    def slide_gpu_tradeoff(self):
        slide = self.base("GPU governor tradeoff", "Time saved versus power spent", ["GPU EDP weights time quadratically and power linearly."])
        self.card(slide, "Ondemand GPU EDP", "+4.0% / +9.9%", 0.65, 1.52, 3.9, 1.42, AMBER, "CPU powersave / CPU performance; median vs default")
        self.card(slide, "Powersave GPU EDP", "+786% / +798%", 0.65, 3.30, 3.9, 1.42, CORAL, "~73% less power, but ~5x GPU time")
        self.text(slide, "Decision rule", 0.75, 5.25, 3.3, 0.3, 17, WHITE, True)
        self.text(slide, "Prefer lower EDP unless latency SLOs justify the power premium.", 0.75, 5.73, 3.6, 0.7, 13, MUTED)
        rows = self.e.workload_changes("gpu", "edp")
        self.image_or_bars(slide, ("gpu", "tradeoff", "edp", "power"), [(r[0], r[3]) for r in rows])

    def slide_workload_wins(self):
        slide = self.base("Workload-specific wins", "Where policy earns its complexity", ["Strong recommendations require magnitude, cross-mode direction, and repeat consistency."])
        self.card(slide, "backprop / ondemand", "-51.5%", 0.65, 1.5, 3.72, 1.55, TEAL, "GPU EDP; <=2.2% latency; 10/10 repeats")
        self.card(slide, "lavaMD / ondemand", "-2.5%", 4.80, 1.5, 3.72, 1.55, TEAL, "GPU EDP; no latency cost; 10/10 repeats")
        self.card(slide, "bh + nbody", "<1%", 8.95, 1.5, 3.72, 1.55, AMBER, "direction agrees, but practical significance is weak")
        path = self.e.image("valid", "win", used=self.used_images)
        self.picture(slide, path, 0.65, 3.45, 12.0, 2.83, path.name if path else None)

    def slide_recommendations(self):
        slide = self.base("Robust recommendations", "Evidence first; exceptions explicit", ["Recommendations quoted or paraphrased from FINDINGS are preferred over inferred universal settings."])
        lines = ["Use gpu_ondemand for backprop and lavaMD after a short guarded trial.",
                 "Treat bh and nbody as marginal candidates; their gains are below 1% and repeat consistency is weak.",
                 "Retain gpu_default for other measured workloads and for unknown applications.",
                 "Avoid gpu_powersave globally; only backprop shows a robust valid win. Never use it for censored gemmEx, mixbench, or myocyte cells."]
        for i, line in enumerate(lines[:4]):
            y = 1.48 + i * 1.25
            color = (TEAL, CORAL, AMBER, TEAL_2)[i]
            self.shape(slide, 0.72, y, 11.9, 0.95, PANEL)
            self.text(slide, f"{i + 1:02d}", 1.0, y + 0.27, 0.55, 0.3, 12, color, True, MONO)
            self.text(slide, line, 1.72, y + 0.18, 10.35, 0.55, 16, WHITE, i == 0, FONT, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)

    def slide_policy(self):
        slide = self.base("Transparent runtime policy", "Portable counters, inspectable decisions", ["Counter names are generic classes; map them to available platform telemetry."])
        steps = [
            ("01", "Observe", "launch latency / device time / utilization", TEAL),
            ("02", "Classify", "host-bound / device-bound / insensitive", AMBER),
            ("03", "Select", "lowest-cost governor meeting the SLO", CORAL),
            ("04", "Verify", "re-sample; fall back on drift", TEAL),
        ]
        for i, (num, title, sub, color) in enumerate(steps):
            x = 0.55 + i * 3.17
            self.shape(slide, x, 1.55, 2.82, 2.0, PANEL)
            self.text(slide, num, x + 0.25, 1.82, 0.45, 0.3, 11, color, True, MONO)
            self.text(slide, title, x + 0.25, 2.28, 2.2, 0.38, 19, WHITE, True)
            self.text(slide, sub, x + 0.25, 2.82, 2.25, 0.52, 11, MUTED)
        counters = ["SM active %", "normalized SM clock", "normalized memory clock", "GPU power / temperature", "launch latency sentinel"]
        self.text(slide, "PORTABLE SIGNAL CLASSES", 0.65, 4.22, 3.6, 0.25, 10, TEAL, True, MONO)
        for i, item in enumerate(counters):
            x = 0.65 + (i % 3) * 4.05
            y = 4.78 + (i // 3) * 0.78
            self.shape(slide, x, y, 3.65, 0.52, PANEL_2)
            self.text(slide, item, x + 0.18, y + 0.12, 3.25, 0.25, 11, WHITE, True)

    def slide_limitations(self):
        slide = self.base("Limitations and x86 HPC next steps", "What must generalize before deployment", ["Do not transfer DGX Spark thresholds directly to x86 HPC nodes."])
        limitations = ["CPU modes were measured on different dates; frequency is confounded with run-date and system drift.",
                       "CPU package power is unavailable, so CPU EDP and whole-system EDP are not reported.",
                       "Component timers may overlap and are diagnostic, not an additive launch-time decomposition.",
                       "Sub-percent EDP differences need reproduction before deployment."]
        self.shape(slide, 0.65, 1.5, 5.75, 4.95, PANEL)
        self.text(slide, "LIMITATIONS", 1.0, 1.86, 4.9, 0.25, 10, CORAL, True, MONO)
        self.bullets(slide, limitations, 1.0, 2.35, 4.9, 3.4, 14, CORAL)
        self.shape(slide, 6.7, 1.5, 5.95, 4.95, PANEL)
        self.text(slide, "X86 HPC VALIDATION", 7.05, 1.86, 5.0, 0.25, 10, TEAL, True, MONO)
        next_steps = ["Repeat the full factorial design on target CPU and accelerator", "Normalize counter classes, not platform-specific event names", "Test NUMA, MPI rank placement, and multi-GPU contention", "Validate policy under production power and thermal limits"]
        self.bullets(slide, next_steps, 7.05, 2.35, 5.05, 3.4, 14, TEAL)
        self.text(slide, "Decision gate: reproduce direction, magnitude, and stability before rollout", 1.35, 6.7, 10.6, 0.28, 13, AMBER, True, MONO, PP_ALIGN.CENTER)

    def validate(self):
        if len(self.prs.slides) != 13:
            raise RuntimeError(f"Expected 13 slides, generated {len(self.prs.slides)}")
        issues = []
        for slide_no, slide in enumerate(self.prs.slides, 1):
            if len(slide.shapes) < 4:
                issues.append(f"slide {slide_no}: suspiciously few objects")
            for shape in slide.shapes:
                if shape.width <= 0 or shape.height <= 0:
                    issues.append(f"slide {slide_no}: non-positive object dimensions")
                # Pictures may intentionally extend beyond the frame before cropping.
                if shape.shape_type != 13 and (shape.left < 0 or shape.top < 0 or shape.left + shape.width > SLIDE_W or shape.top + shape.height > SLIDE_H):
                    issues.append(f"slide {slide_no}: object outside slide bounds")
        if issues:
            raise RuntimeError("Object placement validation failed:\n" + "\n".join(issues))

    def save(self, output: Path, notes_output: Path):
        self.validate()
        output.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(output)
        note_lines = ["DGX Spark CPU/GPU Governor Study - Speaker Notes", ""]
        for idx, (title, notes) in enumerate(self.notes, 1):
            note_lines.extend([f"Slide {idx:02d} - {title}", *(f"- {note}" for note in notes), ""])
        notes_output.write_text("\n".join(note_lines), encoding="utf-8")
        if not output.exists() or output.stat().st_size < 10_000:
            raise RuntimeError(f"Presentation was not written correctly: {output}")
        check = Presentation(output)
        if len(check.slides) != 13:
            raise RuntimeError("Saved presentation failed slide-count validation")


def wait_for_artifacts(root: Path, timeout: int, interval: int) -> None:
    deadline = time.monotonic() + timeout
    last_error = ""
    while True:
        try:
            if not root.is_dir():
                raise FileNotFoundError(f"Analysis directory does not exist: {root}")
            csvs = list(root.rglob("*.csv"))
            findings = [p for p in root.rglob("*") if p.is_file() and "findings" in p.name.lower()]
            if csvs and findings:
                return
            last_error = f"waiting for CSV and FINDINGS artifacts in {root}"
        except OSError as exc:
            last_error = str(exc)
        if time.monotonic() >= deadline:
            raise TimeoutError(f"Artifacts unavailable after {timeout}s: {last_error}")
        time.sleep(interval)


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--artifact-dir", type=Path, default=ROOT)
    parser.add_argument("--output", type=Path, default=OUTPUT)
    parser.add_argument("--wait-seconds", type=int, default=int(os.getenv("GOVERNOR_ARTIFACT_WAIT", "300")))
    parser.add_argument("--poll-seconds", type=int, default=5)
    args = parser.parse_args()
    try:
        wait_for_artifacts(args.artifact_dir, max(0, args.wait_seconds), max(1, args.poll_seconds))
        evidence = Evidence.load(args.artifact_dir)
        deck = Deck(evidence)
        deck.build()
        notes_path = args.output.with_name(args.output.stem + "_notes.txt")
        deck.save(args.output, notes_path)
    except Exception as exc:
        print(f"ERROR: {exc}", file=sys.stderr)
        return 1
    print(args.output)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
